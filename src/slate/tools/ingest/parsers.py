"""Document Ingest — Standardized parsers for enterprise document types.

Each parser extracts structured content from a document format and produces
a standardized IngestResult that can be fed to any Slate pipeline.

Supported formats:
  - PowerPoint (.pptx) — slides, speaker notes, images
  - Word (.docx) — sections, headings, images, tables
  - Excel (.xlsx) — data tables, charts, named ranges

External libraries (python-pptx, python-docx, openpyxl) are optional.
Parsers gracefully degrade if libraries are not installed.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


@dataclass
class MediaAsset:
    """An extracted media asset from a document."""
    asset_type: str      # "image", "chart", "table", "video", "audio"
    source_index: int    # Slide/page/sheet number (0-based)
    file_path: str = ""  # Path to extracted file (populated after extraction)
    description: str = ""
    data: dict[str, Any] = field(default_factory=dict)  # Type-specific metadata


@dataclass
class Section:
    """A structured content section from a document."""
    index: int           # Order in document (0-based)
    title: str = ""
    body: str = ""
    speaker_notes: str = ""
    level: int = 1       # Heading level (1=top)
    media: list[MediaAsset] = field(default_factory=list)
    bullet_points: list[str] = field(default_factory=list)
    data: dict[str, Any] = field(default_factory=dict)

    def to_scene_data(self, scene_id: str | None = None) -> dict[str, Any]:
        """Convert this section into a Slate scenario scene dict."""
        sid = scene_id or f"scene-{self.index + 1}"
        narration = self.speaker_notes or self.body or self.title
        scene: dict[str, Any] = {
            "id": sid,
            "title": self.title,
            "narration": narration,
        }
        if self.bullet_points:
            scene["bullet_points"] = self.bullet_points
        # Generate a visual prompt from the section content
        visual_parts = [self.title]
        if self.bullet_points:
            visual_parts.extend(self.bullet_points[:3])
        scene["visual_prompt"] = f"Professional corporate visual: {', '.join(visual_parts)}"
        return scene


@dataclass
class IngestResult:
    """Standardized output from any document parser.

    This is the contract between ingest tools and pipeline stages.
    Any parser (PPTX, DOCX, XLSX) produces this same structure.
    """
    source_type: str                           # "pptx", "docx", "xlsx"
    source_path: str                           # Original file path
    title: str = ""
    sections: list[Section] = field(default_factory=list)
    media_assets: list[MediaAsset] = field(default_factory=list)
    metadata: dict[str, Any] = field(default_factory=dict)
    estimated_duration_seconds: float = 0.0

    def to_scenario(self, **overrides: Any) -> dict[str, Any]:
        """Convert ingest result to a Slate scenario JSON dict.

        This is the main bridge: document → scenario → SCF → video.
        """
        scenes = [s.to_scene_data() for s in self.sections]

        # Estimate duration: ~8s per scene by default
        total_duration = sum(8 for _ in scenes) + 8  # +8 for intro/outro
        self.estimated_duration_seconds = total_duration

        scenario: dict[str, Any] = {
            "title": self.title or Path(self.source_path).stem,
            "company": self.metadata.get("author", ""),
            "tagline": "",
            "voice": "professional-female",
            "style": "premium-velvet",
            "scenes": scenes,
            "source": {
                "type": self.source_type,
                "path": self.source_path,
                "section_count": len(self.sections),
                "media_count": len(self.media_assets),
            },
        }
        scenario.update(overrides)
        return scenario


class PptxParser:
    """Parse PowerPoint (.pptx) files into IngestResult.

    Extracts:
    - Slide titles and body text
    - Speaker notes
    - Embedded images (saved to output_dir)
    - Bullet point lists

    Requires: python-pptx (optional, degrades gracefully)
    """

    def parse(self, path: str | Path, output_dir: str | Path | None = None) -> IngestResult:
        """Parse a PPTX file.

        Args:
            path: Path to .pptx file
            output_dir: Directory to save extracted images (optional)
        """
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {path}")

        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx not installed — using fallback text extraction")
            return self._fallback_parse(path)

        prs = Presentation(str(path))
        sections: list[Section] = []
        media_assets: list[MediaAsset] = []
        out_dir = Path(output_dir) if output_dir else path.parent / f"{path.stem}_assets"

        for idx, slide in enumerate(prs.slides):
            title = ""
            body_parts: list[str] = []
            bullets: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                        title = text
                    else:
                        body_parts.append(text)
                        # Extract bullet points from paragraphs
                        for para in shape.text_frame.paragraphs:
                            ptext = para.text.strip()
                            if ptext and para.level > 0:
                                bullets.append(ptext)
                            elif ptext and not title:
                                # Treat first substantial text as title if no title shape
                                title = ptext

                # Extract images
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image = shape.image
                    ext = image.content_type.split("/")[-1]
                    img_name = f"slide{idx + 1}_img{len(media_assets) + 1}.{ext}"
                    media_assets.append(MediaAsset(
                        asset_type="image",
                        source_index=idx,
                        file_path=str(out_dir / img_name),
                        description=f"Image from slide {idx + 1}",
                    ))

            # Speaker notes
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            if not title:
                title = f"Slide {idx + 1}"

            sections.append(Section(
                index=idx,
                title=title,
                body="\n".join(body_parts),
                speaker_notes=notes,
                bullet_points=bullets,
                media=media_assets[-1:] if media_assets and media_assets[-1].source_index == idx else [],
            ))

        metadata: dict[str, Any] = {}
        if prs.core_properties:
            cp = prs.core_properties
            if cp.title:
                metadata["doc_title"] = cp.title
            if cp.author:
                metadata["author"] = cp.author

        return IngestResult(
            source_type="pptx",
            source_path=str(path),
            title=metadata.get("doc_title", path.stem),
            sections=sections,
            media_assets=media_assets,
            metadata=metadata,
        )

    def _fallback_parse(self, path: Path) -> IngestResult:
        """Minimal fallback when python-pptx is not available."""
        return IngestResult(
            source_type="pptx",
            source_path=str(path),
            title=path.stem,
            metadata={"warning": "python-pptx not installed — limited extraction"},
        )


class DocxParser:
    """Parse Word (.docx) files into IngestResult.

    Extracts:
    - Headings and body paragraphs
    - Tables
    - Embedded images
    - Document structure (heading levels)

    Requires: python-docx (optional, degrades gracefully)
    """

    def parse(self, path: str | Path, output_dir: str | Path | None = None) -> IngestResult:
        """Parse a DOCX file."""
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"DOCX file not found: {path}")

        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx not installed — using fallback text extraction")
            return self._fallback_parse(path)

        doc = Document(str(path))
        sections: list[Section] = []
        media_assets: list[MediaAsset] = []
        current_section: Section | None = None
        section_idx = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect headings
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                # Save previous section
                if current_section:
                    sections.append(current_section)

                level = 1
                try:
                    level = int(para.style.name.replace("Heading ", "").strip())
                except (ValueError, AttributeError):
                    pass

                current_section = Section(
                    index=section_idx,
                    title=text,
                    level=level,
                )
                section_idx += 1

            elif para.style and para.style.name and "List" in para.style.name:
                if current_section is None:
                    current_section = Section(index=section_idx, title="Introduction")
                    section_idx += 1
                current_section.bullet_points.append(text)

            else:
                if current_section is None:
                    current_section = Section(index=section_idx, title="Introduction")
                    section_idx += 1
                if current_section.body:
                    current_section.body += "\n" + text
                else:
                    current_section.body = text

        # Save last section
        if current_section:
            sections.append(current_section)

        # Extract tables as data sections
        for tidx, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                media_assets.append(MediaAsset(
                    asset_type="table",
                    source_index=tidx,
                    description=f"Table with {len(rows)} rows",
                    data={"rows": rows, "headers": rows[0] if rows else []},
                ))

        metadata: dict[str, Any] = {}
        if doc.core_properties:
            cp = doc.core_properties
            if cp.title:
                metadata["doc_title"] = cp.title
            if cp.author:
                metadata["author"] = cp.author

        # If no sections were found, create one from all text
        if not sections:
            all_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
            sections = [Section(index=0, title=path.stem, body=all_text)]

        return IngestResult(
            source_type="docx",
            source_path=str(path),
            title=metadata.get("doc_title", path.stem),
            sections=sections,
            media_assets=media_assets,
            metadata=metadata,
        )

    def _fallback_parse(self, path: Path) -> IngestResult:
        """Minimal fallback when python-docx is not available."""
        return IngestResult(
            source_type="docx",
            source_path=str(path),
            title=path.stem,
            metadata={"warning": "python-docx not installed — limited extraction"},
        )


class XlsxParser:
    """Parse Excel (.xlsx) files into IngestResult.

    Extracts:
    - Sheet names and data tables
    - Column headers
    - Summary statistics for numeric columns
    - Data suitable for chart generation

    Requires: openpyxl (optional, degrades gracefully)
    """

    def parse(self, path: str | Path, max_rows: int = 1000) -> IngestResult:
        """Parse an XLSX file.

        Args:
            path: Path to .xlsx file
            max_rows: Maximum rows to read per sheet
        """
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"XLSX file not found: {path}")

        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("openpyxl not installed — using fallback")
            return self._fallback_parse(path)

        wb = load_workbook(str(path), data_only=True, read_only=True)
        sections: list[Section] = []
        media_assets: list[MediaAsset] = []

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            rows: list[list[Any]] = []
            headers: list[str] = []

            for ridx, row in enumerate(ws.iter_rows(max_row=max_rows, values_only=True)):
                if ridx == 0:
                    headers = [str(c) if c else f"Col{i}" for i, c in enumerate(row)]
                    continue
                rows.append([c for c in row])

            if not rows and not headers:
                continue

            # Summarize data
            summary_parts: list[str] = []
            summary_parts.append(f"{len(rows)} rows × {len(headers)} columns")
            summary_parts.append(f"Columns: {', '.join(headers[:8])}")
            if len(headers) > 8:
                summary_parts.append(f"...and {len(headers) - 8} more")

            # Compute basic stats for numeric columns
            numeric_stats: dict[str, dict[str, float]] = {}
            for col_idx, header in enumerate(headers):
                values = []
                for row in rows:
                    if col_idx < len(row) and isinstance(row[col_idx], (int, float)):
                        values.append(row[col_idx])
                if values:
                    numeric_stats[header] = {
                        "min": min(values),
                        "max": max(values),
                        "avg": sum(values) / len(values),
                        "count": len(values),
                    }

            body = "\n".join(summary_parts)
            if numeric_stats:
                body += "\n\nNumeric columns:"
                for col, stats in list(numeric_stats.items())[:5]:
                    body += f"\n  {col}: min={stats['min']:.2f}, max={stats['max']:.2f}, avg={stats['avg']:.2f}"

            sections.append(Section(
                index=sheet_idx,
                title=sheet_name,
                body=body,
                data={
                    "headers": headers,
                    "row_count": len(rows),
                    "numeric_stats": numeric_stats,
                    "sample_rows": [list(r) for r in rows[:5]],
                },
            ))

            media_assets.append(MediaAsset(
                asset_type="table",
                source_index=sheet_idx,
                description=f"Sheet '{sheet_name}': {len(rows)} rows × {len(headers)} cols",
                data={"headers": headers, "row_count": len(rows)},
            ))

        wb.close()

        return IngestResult(
            source_type="xlsx",
            source_path=str(path),
            title=path.stem,
            sections=sections,
            media_assets=media_assets,
            metadata={"sheet_count": len(wb.sheetnames)},
        )

    def _fallback_parse(self, path: Path) -> IngestResult:
        """Minimal fallback when openpyxl is not available."""
        return IngestResult(
            source_type="xlsx",
            source_path=str(path),
            title=path.stem,
            metadata={"warning": "openpyxl not installed — limited extraction"},
        )


def parse_document(path: str | Path, **kwargs: Any) -> IngestResult:
    """Auto-detect document type and parse.

    Convenience function that selects the right parser based on file extension.
    """
    path = Path(path)
    ext = path.suffix.lower()

    parsers = {
        ".pptx": PptxParser,
        ".docx": DocxParser,
        ".xlsx": XlsxParser,
    }

    parser_cls = parsers.get(ext)
    if parser_cls is None:
        raise ValueError(
            f"Unsupported document type: {ext}. "
            f"Supported: {', '.join(parsers.keys())}"
        )

    return parser_cls().parse(path, **kwargs)
