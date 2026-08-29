from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "Slate_Talk.pptx"
HERO_IMAGE = OUT_DIR / "assets" / "auditorium-memory.png"
SLIDE5_VIDEO_FRAME = OUT_DIR / "assets" / "slide5-video-frame.png"
SLIDE5_BACKGROUND = OUT_DIR / "assets" / "slide5-vscode-slate.png"
SLIDE5_ARCH_BACKGROUND = OUT_DIR / "assets" / "slide5-vscode-slate-faded.png"
SLIDE7_B15_FRAME = OUT_DIR / "assets" / "slide7-b15-13.png"
CAPABILITY_VIDEO_POSTER = OUT_DIR / "assets" / "slate-capability-video-poster.png"

W = 13.333
H = 7.5

CREAM = "F3EFE7"
INK = "181716"
GOLD = "D6A33D"
GOLD_DARK = "9B6C19"
MUTED = "746F68"
PAPER = "E8DED0"
WHITE = "FFFFFF"
GRAPHITE = "24211F"
TERRACOTTA = "C65C3B"
MOSS = "5F6D59"
BLUE = "506B7A"
OPENMONTAGE = "D14A28"
DISPLAY_FONT = "Georgia"
BODY_FONT = "Gill Sans MT"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def pil_font(size: int, bold=False, mono=False):
    candidates = (
        ["C:/Windows/Fonts/consolab.ttf", "C:/Windows/Fonts/consola.ttf"]
        if mono
        else [
            "C:/Windows/Fonts/seguisb.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
            "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        ]
    )
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def draw_wrapped(draw, text, box, font, fill, spacing=8):
    x, y, width, height = box
    words = text.split()
    lines = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if draw.textbbox((0, 0), candidate, font=font)[2] <= width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    line_height = font.size + spacing
    for index, line in enumerate(lines):
        if y + index * line_height + line_height > y + height:
            break
        draw.text((x, y + index * line_height), line, font=font, fill=fill)


def build_slide5_background() -> Path:
    width, height = 1920, 1080
    chat_width = 458
    top_height = 44
    tabs_height = 42
    terminal_height = 238
    canvas = Image.new("RGB", (width, height), "#181818")
    draw = ImageDraw.Draw(canvas)

    # VS Code shell.
    draw.rectangle((0, 0, width, top_height), fill="#181818")
    draw.rectangle((0, top_height, chat_width, height), fill="#1E1E1E")
    draw.rectangle((chat_width, top_height, width, top_height + tabs_height), fill="#181818")
    draw.line((chat_width, 0, chat_width, height), fill="#343434", width=2)
    draw.line((0, top_height, width, top_height), fill="#303030", width=1)
    draw.text((18, 11), "Visual Studio Code", font=pil_font(18), fill="#C9C9C9")
    draw.text((width - 126, 10), "—   □   ×", font=pil_font(19), fill="#A8A8A8")

    # Chat rail.
    draw.text((22, 60), "CHAT", font=pil_font(16, bold=True), fill="#AFAFAF")
    draw.text((22, 94), "Slate Agent", font=pil_font(24, bold=True), fill="#F0F0F0")
    draw.text((22, 126), "Agentic video production", font=pil_font(17), fill="#858585")
    draw.line((18, 158, chat_width - 18, 158), fill="#343434", width=1)

    def user_bubble(text, y):
        bubble = (118, y, chat_width - 22, y + 88)
        draw.rounded_rectangle(bubble, radius=16, fill="#2B2B2B")
        draw_wrapped(draw, text, (140, y + 18, 274, 60), pil_font(18), "#E7E7E7", 5)

    def slate_message(text, y, height_px=106):
        draw.ellipse((22, y + 2, 48, y + 28), fill="#D6A33D")
        draw.text((31, y + 5), "S", font=pil_font(14, bold=True), fill="#181818")
        draw.text((60, y + 2), "Slate", font=pil_font(18, bold=True), fill="#D6A33D")
        draw_wrapped(draw, text, (60, y + 31, 350, height_px - 34), pil_font(18), "#CCCCCC", 5)

    user_bubble("Create a 60-second launch film from this product brief.", 184)
    slate_message("Script and scene plan ready. Six scenes. Narration, captions and music included.", 292, 128)
    user_bubble("Make the transformation feel more cinematic.", 434)
    slate_message("Art direction revised. Generating the final visuals now.", 542, 104)

    # Production progress block.
    draw.rounded_rectangle((24, 668, chat_width - 24, 788), radius=12, fill="#252526", outline="#3D3D3D")
    draw.text((44, 686), "PRODUCTION", font=pil_font(15, bold=True), fill="#8E8E8E")
    statuses = [("Script", "✓"), ("Scenes", "✓"), ("Assets", "✓"), ("Review", "✓")]
    for index, (label, check) in enumerate(statuses):
        x = 44 + (index % 2) * 190
        y = 722 + (index // 2) * 34
        draw.ellipse((x, y + 4, x + 18, y + 22), fill="#73C991")
        draw.line((x + 4, y + 13, x + 8, y + 18), fill="#181818", width=2)
        draw.line((x + 8, y + 18, x + 15, y + 8), fill="#181818", width=2)
        draw.text((x + 28, y), label, font=pil_font(18), fill="#D5D5D5")

    # Chat composer.
    draw.rounded_rectangle((22, height - 88, chat_width - 22, height - 28), radius=12, fill="#252526", outline="#454545")
    draw.text((42, height - 69), "Ask Slate to revise the cut…", font=pil_font(18), fill="#818181")
    draw.ellipse((chat_width - 62, height - 73, chat_width - 34, height - 45), fill="#007ACC")
    draw.text((chat_width - 54, height - 72), "↑", font=pil_font(19, bold=True), fill="#FFFFFF")

    # Editor tabs and breadcrumb.
    tab_y = top_height
    tabs = ["composition.scf.json", "Slate-launch.mp4", "review-report.md"]
    tab_x = chat_width + 10
    for index, tab in enumerate(tabs):
        tab_w = 214 if index < 2 else 190
        fill = "#1F1F1F" if index == 1 else "#181818"
        draw.rectangle((tab_x, tab_y, tab_x + tab_w, tab_y + tabs_height), fill=fill)
        if index == 1:
            draw.rectangle((tab_x, tab_y, tab_x + tab_w, tab_y + 3), fill="#D6A33D")
        draw.text((tab_x + 18, tab_y + 12), tab, font=pil_font(16), fill="#D0D0D0" if index == 1 else "#8E8E8E")
        tab_x += tab_w

    main_x = chat_width + 16
    main_y = top_height + tabs_height + 40
    main_w = width - main_x - 16
    main_h = height - main_y - terminal_height - 22
    draw.text((main_x, top_height + tabs_height + 10), "projects  ›  slate-launch  ›  renders  ›  final.mp4", font=pil_font(15), fill="#8C8C8C")

    # Video player.
    draw.rounded_rectangle((main_x, main_y, main_x + main_w, main_y + main_h), radius=8, fill="#070A0F")
    frame = Image.open(SLIDE5_VIDEO_FRAME).convert("RGB")
    target_ratio = main_w / (main_h - 58)
    source_ratio = frame.width / frame.height
    if source_ratio > target_ratio:
        crop_w = int(frame.height * target_ratio)
        left = (frame.width - crop_w) // 2
        frame = frame.crop((left, 0, left + crop_w, frame.height))
    else:
        crop_h = int(frame.width / target_ratio)
        top = (frame.height - crop_h) // 2
        frame = frame.crop((0, top, frame.width, top + crop_h))
    frame = frame.resize((main_w, main_h - 58), Image.Resampling.LANCZOS)
    canvas.paste(frame, (main_x, main_y))

    # Video title and player controls.
    overlay = Image.new("RGBA", (main_w, 94), (0, 0, 0, 0))
    overlay_draw = ImageDraw.Draw(overlay)
    for row in range(94):
        alpha = int(185 * (row / 93))
        overlay_draw.rectangle((0, row, main_w, row + 1), fill=(0, 0, 0, alpha))
    canvas.paste(overlay, (main_x, main_y + main_h - 152), overlay)
    draw.text((main_x + 30, main_y + 28), "SLATE  /  FROM INTENT TO FILM", font=pil_font(18, bold=True), fill="#F3E7D2")
    controls_y = main_y + main_h - 42
    draw.polygon(
        [
            (main_x + 24, controls_y - 12),
            (main_x + 24, controls_y + 10),
            (main_x + 40, controls_y - 1),
        ],
        fill="#FFFFFF",
    )
    draw.rectangle((main_x + 60, controls_y, main_x + main_w - 180, controls_y + 4), fill="#727272")
    draw.rectangle((main_x + 60, controls_y, main_x + int(main_w * 0.58), controls_y + 4), fill="#FFFFFF")
    draw.text((main_x + 76, controls_y - 28), "0:42 / 1:00", font=pil_font(15), fill="#E4E4E4")
    draw.text((main_x + main_w - 154, controls_y - 10), "VOL    FULL", font=pil_font(14, bold=True), fill="#E4E4E4")

    # Terminal panel.
    terminal_y = height - terminal_height
    draw.rectangle((chat_width, terminal_y, width, height), fill="#181818")
    draw.line((chat_width, terminal_y, width, terminal_y), fill="#353535", width=2)
    draw.text((main_x, terminal_y + 14), "PROBLEMS    OUTPUT    DEBUG CONSOLE    TERMINAL", font=pil_font(15), fill="#898989")
    draw.rectangle((main_x + 294, terminal_y + 39, main_x + 366, terminal_y + 42), fill="#D6A33D")
    terminal_lines = [
        ("> slate render projects/slate-launch", "#DCDCAA"),
        ("[ok] composition validated", "#73C991"),
        ("[ok] video inspection passed · 8/8", "#73C991"),
        ("output  projects/slate-launch/renders/final.mp4", "#9CDCFE"),
    ]
    for index, (line, color) in enumerate(terminal_lines):
        draw.text((main_x + 8, terminal_y + 62 + index * 34), line, font=pil_font(18, mono=True), fill=color)

    SLIDE5_BACKGROUND.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(SLIDE5_BACKGROUND, quality=95)
    return SLIDE5_BACKGROUND


def build_slide5_architecture_background() -> Path:
    source = Image.open(build_slide5_background()).convert("RGB")
    warm_wash = Image.new("RGB", source.size, "#F3EFE7")
    faded = Image.blend(warm_wash, source, 0.16)
    faded.save(SLIDE5_ARCH_BACKGROUND, quality=95)
    return SLIDE5_ARCH_BACKGROUND


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_chevron(slide, x, y, w=0.30, h=0.42, color=GOLD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=20,
    color=INK,
    font=BODY_FONT,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rule(slide, x, y, w, color=GOLD, height=0.025):
    return add_rect(slide, x, y, w, height, color)


def add_link_text(slide, text, url, x, y, w, h, *, size=20, color=BLUE):
    box = add_text(
        slide,
        text,
        x,
        y,
        w,
        h,
        size=size,
        color=color,
        font=BODY_FONT,
        bold=True,
    )
    box.text_frame.paragraphs[0].runs[0].hyperlink.address = url
    return box


def add_picture_fill(slide, path: Path):
    picture = slide.shapes.add_picture(
        str(path), Inches(0), Inches(0), Inches(W), Inches(H)
    )
    # Source is 3:2; symmetrical vertical crop preserves its composition at 16:9.
    picture.crop_top = 0.078125
    picture.crop_bottom = 0.078125
    return picture


def add_numbered_point(slide, number, title, detail, x, y, width, accent):
    add_text(
        slide,
        number,
        x,
        y,
        0.45,
        0.35,
        size=18,
        color=accent,
        font=BODY_FONT,
        bold=True,
    )
    add_rule(slide, x + 0.58, y + 0.18, width - 0.58, accent, 0.018)
    add_text(
        slide,
        title,
        x,
        y + 0.44,
        width,
        0.42,
        size=24,
        color=INK,
        font=DISPLAY_FONT,
        bold=True,
    )
    add_text(
        slide,
        detail,
        x,
        y + 0.98,
        width,
        1.55,
        size=22,
        color=GRAPHITE,
        font=BODY_FONT,
        line_spacing=1.15,
    )


def add_film_perforations(slide, x, y, count, color):
    for index in range(count):
        add_rect(slide, x, y + index * 0.54, 0.13, 0.24, color, radius=True)


def add_quality(slide, number, label, x, y):
    add_text(
        slide,
        number,
        x,
        y,
        0.38,
        0.25,
        size=9,
        color=GOLD_DARK,
        font="Consolas",
        bold=True,
    )
    add_rule(slide, x + 0.42, y + 0.11, 0.35, GOLD_DARK, 0.018)
    add_text(
        slide,
        label,
        x,
        y + 0.28,
        1.65,
        0.35,
        size=17,
        color=INK,
        font="Aptos Display",
        bold=True,
    )


def add_use_case(slide, index, role, uses, y):
    row_fill = GRAPHITE if index % 2 == 0 else "202023"
    add_rect(slide, 6.42, y, 5.75, 0.84, row_fill)
    add_text(
        slide,
        f"0{index + 1}",
        6.66,
        y + 0.27,
        0.34,
        0.22,
        size=9,
        color=GOLD,
        font="Consolas",
        bold=True,
    )
    add_text(
        slide,
        role,
        7.18,
        y + 0.18,
        1.55,
        0.42,
        size=15,
        color=CREAM,
        font="Aptos Display",
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        uses,
        8.78,
        y + 0.18,
        3.0,
        0.42,
        size=13,
        color=PAPER,
        font="Aptos",
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, number, dark=False):
    color = MUTED if not dark else "77746D"
    add_text(
        slide,
        "SLATE",
        0.65,
        7.10,
        0.8,
        0.18,
        size=8,
        color=color,
        font="Consolas",
        bold=True,
    )
    add_text(
        slide,
        f"0{number}",
        12.15,
        7.10,
        0.5,
        0.18,
        size=8,
        color=color,
        font="Consolas",
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def build_slide_zero(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, INK)

    add_text(slide, "A broader question", 0.70, 0.40, 2.80, 0.34, size=19, color=GOLD, font=BODY_FONT, bold=True)
    add_text(slide, "What else can an agent harness solve?", 0.68, 0.84, 10.30, 0.62, size=35, color=CREAM, font=DISPLAY_FONT)
    add_rect(slide, 10.72, 0.48, 1.88, 0.52, "292623", line="4A4540", radius=True)
    add_text(slide, "BEYOND CODE", 10.82, 0.63, 1.68, 0.22, size=15, color=PAPER, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

    # Central harness: the durable execution environment, not a chat box.
    add_rect(slide, 4.47, 2.36, 4.40, 1.30, "F5F0E7", line=GOLD, radius=True)
    add_text(slide, "CODING HARNESS", 4.79, 2.62, 3.76, 0.30, size=18, color=GOLD_DARK, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Reason + act", 4.79, 3.00, 3.76, 0.34, size=24, color=INK, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "inside a real workspace", 4.79, 3.35, 3.76, 0.22, size=16, color=MUTED, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

    # Five strengths that generalize beyond coding.
    strengths = [
        (0.82, 2.12, 2.78, 0.94, "WORKSPACE", "Files + project context", TERRACOTTA),
        (0.82, 3.53, 2.78, 0.94, "LOCAL TOOLS", "Run real applications", BLUE),
        (9.73, 2.12, 2.78, 0.94, "LIVE WEB", "Research current facts", MOSS),
        (9.73, 3.53, 2.78, 0.94, "SUB-AGENTS", "Delegate specialist work", GOLD),
        (5.17, 4.38, 3.00, 0.94, "HUMAN LOOP", "Review · redirect · approve", TERRACOTTA),
    ]
    for x, y, w, h, label, detail, accent in strengths:
        add_rect(slide, x, y, w, h, "292623", line="4A4540", radius=True)
        add_rect(slide, x, y, 0.09, h, accent)
        add_text(slide, label, x + 0.26, y + 0.17, w - 0.50, 0.22, size=15, color=accent, font=BODY_FONT, bold=True)
        add_text(slide, detail, x + 0.26, y + 0.50, w - 0.50, 0.25, size=18, color=CREAM, font=BODY_FONT, bold=True)

    # Connections suggest one problem-solving substrate.
    add_line(slide, 3.60, 2.59, 4.47, 2.78, "5E5750", 1.5)
    add_line(slide, 3.60, 4.00, 4.47, 3.28, "5E5750", 1.5)
    add_line(slide, 8.87, 2.78, 9.73, 2.59, "5E5750", 1.5)
    add_line(slide, 8.87, 3.28, 9.73, 4.00, "5E5750", 1.5)
    add_line(slide, 6.67, 3.66, 6.67, 4.38, "5E5750", 1.5)

    # Prelude payoff: video production is the example this talk will explore.
    add_rect(slide, 1.10, 5.62, 11.12, 1.22, TERRACOTTA, radius=True)
    add_text(slide, "ONE SUCH PROBLEM", 1.42, 5.78, 10.48, 0.20, size=14, color=WHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Video production", 1.42, 6.02, 10.48, 0.38, size=30, color=WHITE, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "creative  ·  tool-heavy  ·  multi-stage  ·  reviewed", 1.42, 6.48, 10.48, 0.22, size=15, color=WHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, "The rest of this talk is one answer.", 4.62, 7.02, 4.10, 0.24, size=16, color=GOLD, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)


def build_slide_one(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, INK)
    add_picture_fill(slide, HERO_IMAGE)

    add_text(
        slide,
        "Why Slate?",
        0.78,
        0.58,
        2.2,
        0.4,
        size=20,
        color=GOLD,
        font=BODY_FONT,
        bold=True,
    )
    add_text(
        slide,
        "Video stays\nwith people.",
        0.76,
        1.16,
        5.5,
        1.55,
        size=48,
        color=CREAM,
        font=DISPLAY_FONT,
        bold=False,
    )
    add_text(
        slide,
        "Seen   •   Felt   •   Remembered   •   Shared",
        0.80,
        3.02,
        5.6,
        0.48,
        size=23,
        color=PAPER,
        font=BODY_FONT,
    )

    add_rect(slide, 0, 5.88, W, 1.62, INK)
    use_cases = [
        ("PMs", "Launches · roadmaps"),
        ("Engineers", "Architecture · demos"),
        ("Trainers", "Onboarding · readiness"),
        ("Leaders", "Strategy · updates"),
    ]
    for index, (role, uses) in enumerate(use_cases):
        x = 0.78 + index * 3.13
        add_text(
            slide,
            role,
            x,
            6.18,
            2.7,
            0.38,
            size=23,
            color=CREAM,
            font=DISPLAY_FONT,
            bold=True,
        )
        add_text(
            slide,
            uses,
            x,
            6.66,
            2.72,
            0.34,
            size=19,
            color=GOLD if index in {0, 3} else PAPER,
            font=BODY_FONT,
        )


def add_discipline(slide, index, label, x, y, w=2.34):
    fill = CREAM if index in {0, 3, 6, 9} else GRAPHITE
    text_color = INK if fill == CREAM else CREAM
    line_color = None if fill == CREAM else "363538"
    add_rect(slide, x, y, w, 0.7, fill, line=line_color)
    add_text(
        slide,
        f"{index + 1:02d}",
        x + 0.16,
        y + 0.22,
        0.3,
        0.2,
        size=8,
        color=GOLD_DARK if fill == CREAM else GOLD,
        font="Consolas",
        bold=True,
    )
    add_text(
        slide,
        label,
        x + 0.53,
        y + 0.16,
        w - 0.66,
        0.34,
        size=11,
        color=text_color,
        font="Aptos Display",
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build_slide_two(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)

    add_text(
        slide,
        "The barrier",
        0.78,
        0.54,
        2.0,
        0.38,
        size=20,
        color=TERRACOTTA,
        font=BODY_FONT,
        bold=True,
    )
    add_text(
        slide,
        "One video.\nTen crafts.",
        0.76,
        1.04,
        4.2,
        1.32,
        size=43,
        color=INK,
        font=DISPLAY_FONT,
        bold=False,
    )
    add_numbered_point(slide, "01", "Story", "Script\nScene plan", 5.18, 0.72, 2.1, TERRACOTTA)
    add_numbered_point(
        slide,
        "02",
        "Craft",
        "Narration\nImages + clips\nMotion graphics\nAnimation\nTransitions",
        7.72,
        0.72,
        2.25,
        MOSS,
    )
    add_numbered_point(
        slide,
        "03",
        "Finish",
        "Music + captions\nEdit + render\nWeb + mobile cuts",
        10.50,
        0.72,
        2.15,
        BLUE,
    )

    add_rect(slide, 0, 5.47, W, 2.03, INK)
    add_text(slide, "You", 0.78, 5.80, 1.75, 0.48, size=30, color=CREAM, font=DISPLAY_FONT, bold=True)
    add_text(slide, "Creative head", 0.78, 6.36, 2.2, 0.4, size=22, color=GOLD, font=BODY_FONT)
    add_text(slide, "→", 3.20, 6.02, 0.55, 0.44, size=28, color=TERRACOTTA, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "GitHub Copilot", 4.15, 5.80, 3.45, 0.48, size=30, color=CREAM, font=DISPLAY_FONT, bold=True)
    add_text(slide, "Director", 4.15, 6.36, 1.75, 0.4, size=22, color=GOLD, font=BODY_FONT)
    add_text(slide, "→", 8.04, 6.02, 0.55, 0.44, size=28, color=MOSS, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Finished video", 9.06, 5.80, 3.25, 0.48, size=30, color=CREAM, font=DISPLAY_FONT, bold=True)
    add_text(slide, "Web · mobile · stage", 9.06, 6.36, 3.45, 0.4, size=22, color=GOLD, font=BODY_FONT)


def build_slide_three(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, INK)
    slide.shapes.add_picture(
        str(CAPABILITY_VIDEO_POSTER),
        Inches(0),
        Inches(0),
        Inches(W),
        Inches(H),
    )


def build_slide_four(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)

    add_text(
        slide,
        "How Slate came about",
        0.78,
        0.42,
        3.6,
        0.42,
        size=20,
        color=TERRACOTTA,
        font=BODY_FONT,
        bold=True,
    )
    add_text(
        slide,
        "Open Source Project. Same creator.",
        0.76,
        0.88,
        9.0,
        0.62,
        size=36,
        color=INK,
        font=DISPLAY_FONT,
    )

    # Two equal-height panels create one continuous origin story.
    add_rect(slide, 0.76, 1.76, 5.52, 5.08, INK)
    add_rect(slide, 0.76, 1.76, 0.13, 5.08, OPENMONTAGE)
    add_rect(slide, 6.54, 1.76, 6.03, 5.08, "FBF8F1", line="D8D0C4")

    add_text(
        slide,
        "Inspired by",
        1.17,
        2.10,
        2.1,
        0.30,
        size=18,
        color=OPENMONTAGE,
        font=BODY_FONT,
        bold=True,
    )
    add_text(
        slide,
        "OpenMontage",
        1.17,
        2.44,
        3.55,
        0.58,
        size=31,
        color=CREAM,
        font=DISPLAY_FONT,
        bold=True,
    )
    add_text(
        slide,
        "World’s first open-source\nagentic video production system",
        1.17,
        3.15,
        4.28,
        0.88,
        size=22,
        color=PAPER,
        font=BODY_FONT,
        line_spacing=1.08,
    )
    add_text(slide, "50K+", 1.17, 4.28, 1.55, 0.58, size=32, color=GOLD, font=DISPLAY_FONT, bold=True)
    add_text(slide, "GitHub stars", 1.17, 4.86, 1.65, 0.34, size=18, color=CREAM, font=BODY_FONT)
    add_text(slide, "6.3K", 3.20, 4.28, 1.35, 0.58, size=32, color=GOLD, font=DISPLAY_FONT, bold=True)
    add_text(slide, "forks", 3.20, 4.86, 1.25, 0.34, size=18, color=CREAM, font=BODY_FONT)
    add_text(slide, "Viral across social media", 1.17, 5.48, 3.20, 0.34, size=21, color=PAPER, font=BODY_FONT, bold=True)
    add_link_text(
        slide,
        "github.com/calesthio/OpenMontage  ↗",
        "https://github.com/calesthio/OpenMontage",
        1.17,
        6.15,
        4.45,
        0.34,
        size=18,
        color=BLUE,
    )

    add_text(
        slide,
        "The question",
        6.92,
        2.10,
        2.0,
        0.36,
        size=20,
        color=MOSS,
        font=BODY_FONT,
        bold=True,
    )
    add_text(
        slide,
        "Why stop at creators?",
        6.90,
        2.52,
        5.20,
        0.72,
        size=36,
        color=INK,
        font=DISPLAY_FONT,
    )
    add_rule(slide, 6.92, 3.40, 5.25, TERRACOTTA, 0.035)

    # Aligned evidence tiles replace the earlier floating equation.
    add_rect(slide, 6.92, 3.75, 2.18, 0.96, "F4E0D8")
    add_rect(slide, 9.72, 3.75, 2.43, 0.96, "E2E7DE")
    add_text(slide, "Real-world\nneed", 7.05, 3.84, 1.92, 0.70, size=19, color=GRAPHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.9)
    add_text(slide, "Open-source\nsignal", 9.86, 3.84, 2.16, 0.70, size=19, color=GRAPHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.9)
    add_text(slide, "+", 9.25, 3.99, 0.34, 0.38, size=25, color=TERRACOTTA, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.CENTER)

    add_text(
        slide,
        "PMs  ·  Engineers  ·  Trainers  ·  Leaders",
        6.92,
        4.86,
        5.7,
        0.42,
        size=21,
        color=MUTED,
        font=BODY_FONT,
    )
    add_text(
        slide,
        "Slate",
        6.90,
        5.39,
        2.45,
        0.66,
        size=40,
        color=TERRACOTTA,
        font=DISPLAY_FONT,
        bold=True,
    )
    add_text(
        slide,
        "Agentic video for business",
        6.92,
        6.04,
        3.65,
        0.44,
        size=24,
        color=INK,
        font=BODY_FONT,
        bold=True,
    )


def add_stage_node(slide, label, x, y, w, fill, *, color=INK, line=None):
    add_rect(slide, x, y, w, 0.64, fill, line=line, radius=True)
    add_text(
        slide,
        label,
        x + 0.08,
        y + 0.13,
        w - 0.16,
        0.34,
        size=18,
        color=color,
        font=BODY_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_capability_row(slide, label, model, x, y, w, accent):
    add_rect(slide, x, y, w, 0.44, "282523", line="403B37", radius=True)
    add_rect(slide, x, y, 0.08, 0.44, accent)
    add_text(slide, label, x + 0.18, y + 0.10, 1.10, 0.24, size=16, color=CREAM, font=BODY_FONT, bold=True)
    add_text(slide, model, x + 1.25, y + 0.10, w - 1.38, 0.24, size=16, color=PAPER, font=BODY_FONT, align=PP_ALIGN.RIGHT)


def build_slide_five(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = build_slide5_architecture_background()
    slide.shapes.add_picture(str(background), Inches(0), Inches(0), Inches(W), Inches(H))

    add_text(slide, "High-level architecture", 0.68, 0.34, 3.4, 0.36, size=19, color=TERRACOTTA, font=BODY_FONT, bold=True)
    add_text(slide, "From an idea to a reviewed video", 0.66, 0.74, 7.2, 0.58, size=36, color=INK, font=DISPLAY_FONT)

    flow = [
        ("Direct", "You + Copilot", TERRACOTTA),
        ("Plan", "Research · script\nScene plan", GOLD_DARK),
        ("Create", "Images · video\nVoice · captions", MOSS),
        ("Assemble", "Motion · music\nEdit", BLUE),
        ("Inspect", "Quality review", GOLD),
        ("Deliver", "Video", TERRACOTTA),
    ]
    x = 0.68
    widths = [1.50, 1.74, 1.74, 1.74, 1.58, 1.36]
    for index, ((title, detail, accent), width) in enumerate(zip(flow, widths)):
        add_rect(slide, x, 2.02, width, 1.26, "FBF8F1", line=accent, radius=True)
        add_text(slide, title, x + 0.08, 2.25, width - 0.16, 0.34, size=22, color=INK, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, detail, x + 0.10, 2.67, width - 0.20, 0.48, size=15, color=MUTED, font=BODY_FONT, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        if index < len(flow) - 1:
            add_chevron(slide, x + width + 0.08, 2.43, 0.28, 0.42, accent)
        x += width + 0.46

    add_rect(slide, 2.30, 4.08, 2.92, 1.34, INK, radius=True)
    add_text(slide, "Azure AI Foundry", 2.53, 4.29, 2.46, 0.34, size=21, color=CREAM, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Images · video\nNarration · captions", 2.53, 4.70, 2.46, 0.52, size=16, color=PAPER, font=BODY_FONT, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.92)

    add_rect(slide, 7.50, 4.08, 3.42, 1.34, GRAPHITE, radius=True)
    add_text(slide, "Inspection + review", 7.76, 4.29, 2.90, 0.34, size=21, color=CREAM, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Video · audio\nCaptions · brand", 7.76, 4.70, 2.90, 0.52, size=16, color=PAPER, font=BODY_FONT, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.92)

    add_line(slide, 3.76, 4.08, 3.76, 3.61, TERRACOTTA, 2)
    add_line(slide, 3.76, 3.61, 5.60, 3.61, TERRACOTTA, 2)
    add_line(slide, 5.60, 3.61, 5.60, 3.28, TERRACOTTA, 2)
    add_chevron(slide, 5.46, 3.18, 0.28, 0.34, TERRACOTTA)

    add_line(slide, 9.21, 4.08, 9.21, 3.28, GOLD, 2)
    add_chevron(slide, 9.07, 3.18, 0.28, 0.34, GOLD)

    add_line(slide, 9.21, 5.42, 9.21, 6.08, TERRACOTTA, 2)
    add_line(slide, 9.21, 6.08, 4.02, 6.08, TERRACOTTA, 2)
    add_line(slide, 4.02, 6.08, 4.02, 5.74, TERRACOTTA, 2)
    add_text(slide, "revise", 6.20, 5.86, 0.78, 0.28, size=16, color=TERRACOTTA, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 1.15, 6.48, 11.02, 0.54, "E4DDD2", radius=True)
    add_text(slide, "Human approvals", 1.50, 6.63, 1.80, 0.24, size=16, color=GRAPHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Brand", 4.22, 6.63, 1.00, 0.24, size=16, color=GRAPHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Cost controls", 6.34, 6.63, 1.52, 0.24, size=16, color=GRAPHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Audit trail", 9.20, 6.63, 1.38, 0.24, size=16, color=GRAPHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)


def build_walkthrough_cue(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, INK)
    add_text(
        slide,
        "Walkthrough",
        2.00,
        3.08,
        9.33,
        0.86,
        size=44,
        color=CREAM,
        font=DISPLAY_FONT,
        bold=False,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def build_slide_six(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, INK)

    add_text(slide, "The operating model", 0.68, 0.36, 3.1, 0.34, size=19, color=GOLD, font=BODY_FONT, bold=True)
    add_text(slide, "Creative freedom. Production discipline.", 0.66, 0.80, 7.76, 0.54, size=31, color=CREAM, font=DISPLAY_FONT)

    # Capability is easy; dependable delivery is the hard part.
    add_rect(slide, 8.70, 0.34, 1.48, 0.92, "2A2826", line="4A4641", radius=True)
    add_text(slide, "EASY", 8.91, 0.50, 1.06, 0.26, size=17, color=GOLD, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Tools + skills", 8.84, 0.83, 1.20, 0.24, size=16, color=CREAM, font=BODY_FONT, align=PP_ALIGN.CENTER)
    add_chevron(slide, 10.36, 0.60, 0.34, 0.42, TERRACOTTA)
    add_rect(slide, 10.88, 0.34, 1.80, 0.92, TERRACOTTA, radius=True)
    add_text(slide, "HARD", 11.08, 0.50, 1.40, 0.26, size=17, color=WHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "End-to-end video", 11.00, 0.83, 1.56, 0.24, size=16, color=WHITE, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Two aligned responsibility planes.
    add_rect(slide, 0.68, 1.70, 5.62, 4.64, "24211F", line="4A4540", radius=True)
    add_rect(slide, 7.03, 1.70, 5.62, 4.64, "F5F0E7", line="D7CAB8", radius=True)

    add_text(slide, "Agent creativity", 1.02, 2.04, 3.50, 0.42, size=25, color=CREAM, font=DISPLAY_FONT, bold=True)
    add_text(slide, "Deterministic production", 7.38, 2.04, 4.58, 0.42, size=25, color=INK, font=DISPLAY_FONT, bold=True)

    creative = [
        ("01", "Meaning", "What should this say?"),
        ("02", "Emotion", "What should it feel like?"),
        ("03", "Story", "How should it unfold?"),
        ("04", "Visual world", "How should it look?"),
        ("05", "Treatment", "Which idea is strongest?"),
    ]
    production = [
        ("Facts verified", "Sources captured"),
        ("Direction approved", "Human checkpoint"),
        ("Assets present", "Nothing imaginary"),
        ("Timing valid", "Audio fits the scene"),
        ("Video inspected", "Pass before delivery"),
    ]

    for index, ((number, title, question), (check_title, check_detail)) in enumerate(zip(creative, production)):
        y = 2.66 + index * 0.68

        # Creative side: numbered thought cards.
        add_rect(slide, 1.02, y, 4.92, 0.54, "2F2B28", line="49433E", radius=True)
        add_text(slide, number, 1.18, y + 0.13, 0.42, 0.24, size=16, color=GOLD, font=BODY_FONT, bold=True)
        add_text(slide, title, 1.72, y + 0.10, 1.28, 0.26, size=19, color=CREAM, font=DISPLAY_FONT, bold=True)
        add_text(slide, question, 3.02, y + 0.12, 2.64, 0.24, size=16, color=PAPER, font=BODY_FONT, align=PP_ALIGN.RIGHT)

        # Production side: explicit pass conditions.
        add_rect(slide, 7.38, y, 4.92, 0.54, WHITE, line="DDD3C6", radius=True)
        cx = 7.56
        cy = y + 0.15
        draw_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(0.23), Inches(0.23))
        draw_circle.fill.solid()
        draw_circle.fill.fore_color.rgb = rgb(MOSS)
        draw_circle.line.fill.background()
        add_text(slide, "✓", cx + 0.01, cy - 0.01, 0.21, 0.20, size=13, color=WHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, check_title, 7.94, y + 0.10, 2.30, 0.26, size=18, color=INK, font=BODY_FONT, bold=True)
        add_text(slide, check_detail, 10.25, y + 0.12, 1.81, 0.24, size=16, color=MUTED, font=BODY_FONT, align=PP_ALIGN.RIGHT)

        # Pairing line through the Slate spine.
        add_line(slide, 6.30, y + 0.27, 7.03, y + 0.27, "5E5750", 1.2)

    # Slate is the conversion point between free judgment and repeatable production.
    add_rect(slide, 6.10, 2.36, 0.43, 3.66, TERRACOTTA, radius=True)
    for index, letter in enumerate("SLATE"):
        add_text(slide, letter, 6.12, 2.67 + index * 0.57, 0.39, 0.28, size=17, color=WHITE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 1.54, 6.62, 10.23, 0.48, "302C29", line="4A4540", radius=True)
    add_text(slide, "No reliance on agent memory", 1.84, 6.74, 3.10, 0.24, size=17, color=GOLD, font=BODY_FONT, bold=True)
    add_text(slide, "Choices become visible  ·  approved  ·  checked  ·  repeatable", 4.66, 6.74, 6.72, 0.24, size=17, color=CREAM, font=BODY_FONT, align=PP_ALIGN.RIGHT)


def build_slide_seven(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)

    add_text(slide, "How durability works", 0.68, 0.36, 3.20, 0.34, size=19, color=TERRACOTTA, font=BODY_FONT, bold=True)
    add_text(slide, "The plan survives the conversation", 0.66, 0.80, 7.8, 0.58, size=36, color=INK, font=DISPLAY_FONT)
    add_text(slide, "One scene, three layers of evidence", 9.08, 0.88, 3.54, 0.30, size=18, color=MUTED, font=BODY_FONT, bold=True, align=PP_ALIGN.RIGHT)

    # Step headings.
    headings = [
        (0.68, "01", "Preserved artifacts", TERRACOTTA),
        (4.38, "02", "Scene blueprint", GOLD_DARK),
        (7.78, "03", "Rendered evidence", MOSS),
    ]
    for x, number, title, accent in headings:
        add_text(slide, number, x, 1.54, 0.42, 0.28, size=17, color=accent, font=BODY_FONT, bold=True)
        add_rule(slide, x + 0.50, 1.68, 0.50, accent, 0.025)
        add_text(slide, title, x, 1.94, 3.10, 0.40, size=23, color=INK, font=DISPLAY_FONT, bold=True)

    # 1. A tactile dossier of real production snippets.
    add_rect(slide, 0.68, 2.48, 3.28, 3.75, "EEE7DB", line="D7CEC1", radius=True)
    snippets = [
        ("BRIEF", "Engineering leaders"),
        ("APPROVED SCRIPT", "Verify before you connect."),
        ("SCENE PLAN", "Genuine green · Fake red"),
        ("DECISION", "Treatment approved"),
    ]
    for index, (source, value) in enumerate(snippets):
        y = 2.76 + index * 0.78
        card_fill = "FFFDF9" if index % 2 == 0 else "F8F3EA"
        add_rect(slide, 0.94, y, 2.76, 0.64, card_fill, line="D9D0C3", radius=True)
        add_text(slide, source, 1.10, y + 0.09, 2.36, 0.18, size=13, color=TERRACOTTA if index < 3 else MOSS, font=BODY_FONT, bold=True)
        add_text(slide, value, 1.10, y + 0.32, 2.36, 0.22, size=15, color=INK, font=BODY_FONT, bold=True)

    # 2. Plain-language executable scene blueprint.
    add_rect(slide, 4.38, 2.48, 2.92, 3.75, INK, radius=True)
    add_text(slide, "SCENE 04", 4.72, 2.78, 2.24, 0.34, size=25, color=GOLD, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "PRODUCTION BLUEPRINT", 4.72, 3.17, 2.24, 0.20, size=13, color=PAPER, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    blueprint = [
        ("Window", "12.6 sec"),
        ("Voice", "12.1 sec"),
        ("Visual", "Trust reveal"),
        ("Motion", "5 beats"),
        ("Next", "Soft fade"),
        ("Media", "All found"),
    ]
    for index, (label, value) in enumerate(blueprint):
        y = 3.55 + index * 0.39
        add_text(slide, label, 4.72, y, 0.82, 0.22, size=15, color=MUTED, font=BODY_FONT)
        add_text(slide, value, 5.56, y, 1.40, 0.22, size=16, color=CREAM, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.RIGHT)
    add_rule(slide, 4.72, 5.88, 2.24, GOLD, 0.025)
    add_text(slide, "READY TO RENDER", 4.72, 5.96, 2.24, 0.20, size=13, color=GOLD, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

    # Hand-off arrows.
    add_chevron(slide, 4.04, 4.00, 0.27, 0.46, TERRACOTTA)
    add_chevron(slide, 7.43, 4.00, 0.27, 0.46, GOLD)

    # 3. Real rendered scene with review stamps.
    add_rect(slide, 7.78, 2.48, 4.86, 3.75, "0B0D12", line="34312F", radius=True)
    picture = slide.shapes.add_picture(
        str(SLIDE7_B15_FRAME),
        Inches(7.94), Inches(2.64), Inches(4.54), Inches(2.56),
    )
    picture.crop_bottom = 0.10
    add_rect(slide, 7.94, 5.27, 4.54, 0.70, "23211F", radius=True)
    stamps = [
        ("TIMING", "PASS", 8.15, GOLD),
        ("MEDIA", "PASS", 9.65, BLUE),
        ("INSPECTION", "PASS", 11.05, "73C991"),
    ]
    for label, value, x, accent in stamps:
        add_text(slide, label, x, 5.38, 1.14, 0.18, size=12, color=MUTED, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, value, x, 5.61, 1.14, 0.22, size=17, color=accent, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)

    # The durable consequence.
    add_rect(slide, 0.98, 6.58, 11.36, 0.46, "E4DDD2", radius=True)
    add_text(slide, "Resume", 1.46, 6.70, 1.08, 0.22, size=17, color=TERRACOTTA, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Revise", 3.18, 6.70, 1.08, 0.22, size=17, color=GOLD_DARK, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Reproduce", 4.96, 6.70, 1.36, 0.22, size=17, color=BLUE, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "without the original agent or chat history", 7.02, 6.70, 4.84, 0.22, size=17, color=INK, font=DISPLAY_FONT, bold=True, align=PP_ALIGN.RIGHT)


def build_slide_eight(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)

    add_text(slide, "SLATE TODAY", 0.68, 0.38, 2.10, 0.28, size=16, color=TERRACOTTA, font=BODY_FONT, bold=True)
    add_text(slide, "Current Operating Realities", 0.66, 0.82, 8.10, 0.62, size=37, color=INK, font=DISPLAY_FONT)
    add_rule(slide, 0.68, 1.58, 11.96, GRAPHITE, 0.018)

    columns = [
        (
            "01",
            "A project,\nnot yet a product",
            "Built for exploration and extension, not turnkey adoption.",
            TERRACOTTA,
        ),
        (
            "02",
            "Azure infrastructure\nrequired",
            "Foundry resources and model deployments power generated media and speech.",
            GOLD_DARK,
        ),
        (
            "03",
            "Production\ntakes time",
            "Advanced models, video generation, and rendering can take several minutes.",
            BLUE,
        ),
    ]
    starts = [0.68, 4.50, 8.32]

    for index, ((number, title, detail, accent), x) in enumerate(zip(columns, starts)):
        if index:
            add_rule(slide, x - 0.36, 1.94, 0.018, "CFC6B9", 4.58)

        add_text(slide, number, x, 1.98, 0.72, 0.40, size=20, color=accent, font=BODY_FONT, bold=True)
        add_rule(slide, x + 0.78, 2.18, 2.54, accent, 0.025)

        # A simple editable visual metaphor for each operating reality.
        if index == 0:
            add_rect(slide, x + 0.18, 2.66, 2.62, 1.44, "E6DDD0", line="C9BFB1", radius=True)
            add_rect(slide, x + 0.34, 2.50, 2.62, 1.44, "FBF8F1", line=accent, radius=True)
            add_text(slide, "PROJECT", x + 0.62, 2.80, 1.55, 0.24, size=15, color=accent, font=BODY_FONT, bold=True)
            add_rule(slide, x + 0.62, 3.18, 1.78, GRAPHITE, 0.025)
            add_rule(slide, x + 0.62, 3.47, 1.28, "BEB4A7", 0.025)
            add_rect(slide, x + 2.48, 3.44, 0.28, 0.28, accent, radius=True)
        elif index == 1:
            add_rect(slide, x + 0.26, 2.56, 2.86, 1.46, INK, radius=True)
            add_text(slide, "AZURE AI FOUNDRY", x + 0.52, 2.80, 2.34, 0.25, size=15, color=CREAM, font=BODY_FONT, bold=True, align=PP_ALIGN.CENTER)
            node_xs = [x + 0.64, x + 1.42, x + 2.20]
            node_colors = [TERRACOTTA, GOLD, BLUE]
            for node_x, node_color in zip(node_xs, node_colors):
                node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(node_x), Inches(3.30), Inches(0.42), Inches(0.42))
                node.fill.solid()
                node.fill.fore_color.rgb = rgb(node_color)
                node.line.fill.background()
            add_line(slide, x + 1.06, 3.51, x + 1.42, 3.51, PAPER, 1.5)
            add_line(slide, x + 1.84, 3.51, x + 2.20, 3.51, PAPER, 1.5)
        else:
            clock = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.78), Inches(2.46), Inches(1.72), Inches(1.72))
            clock.fill.solid()
            clock.fill.fore_color.rgb = rgb("FBF8F1")
            clock.line.color.rgb = rgb(accent)
            clock.line.width = Pt(3)
            add_line(slide, x + 1.64, 3.32, x + 1.64, 2.78, accent, 3)
            add_line(slide, x + 1.64, 3.32, x + 2.04, 3.53, accent, 3)
            add_rect(slide, x + 2.82, 2.72, 0.42, 0.22, TERRACOTTA, radius=True)
            add_rect(slide, x + 2.82, 3.10, 0.66, 0.22, GOLD, radius=True)
            add_rect(slide, x + 2.82, 3.48, 0.90, 0.22, BLUE, radius=True)

        title_box = add_text(slide, title, x, 4.50, 3.22, 0.94, size=23, color=INK, font=DISPLAY_FONT, bold=True, line_spacing=0.96)
        title_box.text_frame.word_wrap = True
        detail_box = add_text(slide, detail, x, 5.64, 3.12, 0.92, size=16, color=GRAPHITE, font=BODY_FONT, line_spacing=1.05)
        detail_box.text_frame.word_wrap = True

    add_text(slide, "Clear expectations make the experiment stronger.", 0.68, 6.95, 6.70, 0.28, size=17, color=MUTED, font=BODY_FONT, bold=True)
    add_text(slide, "09", 12.00, 6.95, 0.62, 0.26, size=15, color=MUTED, font="Consolas", bold=True, align=PP_ALIGN.RIGHT)


def build_slide_nine(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, INK)

    repository_url = "https://github.com/gim-home/SlateStudio"
    videos_url = "https://onedrive.cloud.microsoft/my?id=%2Fa%40j4kj6sq6%2FDocuments%2FSlate%20Videos&share=cgrnEJQhbtxNTIlwUCxrF4lkEgUCPRI86l%2DYAsXNHPOzREEa4w&CT=1787800500867&OR=OWA%2DNT%2DMail&SI=SentItems&SLSync=Y"
    channel_url = "https://teams.microsoft.com/l/team/19%3AVAkxNGrvPQew7Y125vVQH2qWb5emtU1PAxGdm0P3IQA1%40thread.tacv2/conversations?groupId=2cddd44a-d6c1-41e3-8da3-a6e9a1cd1669&tenantId=72f988bf-86f1-41af-91ab-2d7cd011db47"
    office_hours_url = "https://teams.microsoft.com/l/meetup-join/19%3ameeting_OThjODE3OTktMGJkZS00MDczLTk5YjItNDBiMzNiMzkxOTg1%40thread.v2/0?context=%7b%22Tid%22%3a%2272f988bf-86f1-41af-91ab-2d7cd011db47%22%2c%22Oid%22%3a%22d59af512-38f6-4b44-a63d-57d11e8b6fe7%22%7d"

    add_text(slide, "NEXT STEPS", 0.68, 0.38, 1.90, 0.28, size=16, color=GOLD, font=BODY_FONT, bold=True)
    add_text(slide, "Get started today", 0.66, 0.80, 7.30, 0.62, size=38, color=CREAM, font=DISPLAY_FONT)
    add_text(slide, "Four ways into the project", 9.02, 0.92, 3.60, 0.28, size=18, color=PAPER, font=BODY_FONT, bold=True, align=PP_ALIGN.RIGHT)
    add_rule(slide, 0.68, 1.56, 11.96, "4A4540", 0.018)

    actions = [
        ("01", "START", "Clone SlateStudio", "Point Copilot or your coding agent to the repository and ask it to set up the workspace.", "Open repository  →", repository_url, TERRACOTTA),
        ("02", "WATCH", "Some reference videos made using Slate", "See complete examples created through the Slate production workflow.", "View videos  →", videos_url, GOLD),
        ("03", "JOIN", "General | Slate", "Connect with the Slate community for questions, updates, and production notes.", "Open Teams channel  →", channel_url, MOSS),
        ("04", "OFFICE HOURS", "Every Wednesday · 11:00 AM PST", "Bring a project, a production question, or simply watch Slate work.", "Join office hours  →", office_hours_url, BLUE),
    ]

    for index, (number, kicker, title, detail, link_label, url, accent) in enumerate(actions):
        y = 1.88 + index * 1.22
        row_fill = "24211F" if index % 2 == 0 else "211F1D"
        add_rect(slide, 0.68, y, 11.96, 1.00, row_fill, line="46413C", radius=True)
        add_rect(slide, 0.68, y, 0.08, 1.00, accent)
        add_text(slide, number, 0.96, y + 0.33, 0.48, 0.25, size=15, color=accent, font="Consolas", bold=True)
        add_text(slide, kicker, 1.52, y + 0.17, 1.36, 0.20, size=13, color=accent, font=BODY_FONT, bold=True)
        title_box = add_text(slide, title, 1.52, y + 0.40, 3.78, 0.46, size=18, color=CREAM, font=DISPLAY_FONT, bold=True, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        title_box.text_frame.word_wrap = True
        detail_box = add_text(slide, detail, 5.52, y + 0.22, 4.02, 0.56, size=14, color=PAPER, font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.96)
        detail_box.text_frame.word_wrap = True
        link_box = add_text(slide, link_label, 9.82, y + 0.35, 2.44, 0.30, size=14, color=accent, font=BODY_FONT, bold=True, align=PP_ALIGN.RIGHT)
        link_box.text_frame.word_wrap = True
        link_box.text_frame.paragraphs[0].runs[0].font.underline = True
        link_box.click_action.hyperlink.address = url

    add_text(slide, "SLATE", 0.68, 7.06, 0.72, 0.18, size=10, color="77746D", font="Consolas", bold=True)
    add_text(slide, "11", 12.00, 7.04, 0.62, 0.22, size=13, color="77746D", font="Consolas", bold=True, align=PP_ALIGN.RIGHT)


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "Why Slate?"
    prs.core_properties.subject = "Slate talk opening"
    prs.core_properties.author = "Slate"
    prs.core_properties.keywords = "Slate, GitHub Copilot, video production"

    build_slide_zero(prs)
    build_slide_one(prs)
    build_slide_two(prs)
    build_slide_three(prs)
    build_slide_four(prs)
    build_walkthrough_cue(prs)
    build_slide_five(prs)
    build_slide_six(prs)
    build_slide_seven(prs)
    build_slide_eight(prs)
    build_slide_nine(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    return OUT_FILE


if __name__ == "__main__":
    print(build_deck())